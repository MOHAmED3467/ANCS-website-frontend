import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt

base_dir = os.path.dirname(__file__)
assets_dir = os.path.normpath(os.path.join(base_dir, '..', 'src', 'assets'))
output_path = os.path.normpath(os.path.join(base_dir, '..', 'public', 'ANCS-presentation.pptx'))

presentation = Presentation()


def set_slide_background_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


def add_title(slide, title_text):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1.0))
    title_tf = title_box.text_frame
    title_tf.text = title_text
    title_tf.paragraphs[0].font.name = 'Calibri'
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_tf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT


def add_paragraphs(slide, lines, left=Inches(0.7), top=Inches(1.7), width=Inches(8.6), height=Inches(4.0)):
    body = slide.shapes.add_textbox(left, top, width, height).text_frame
    body.word_wrap = True
    body.margin_bottom = Pt(8)
    body.margin_top = Pt(8)
    body.margin_left = Pt(8)
    body.margin_right = Pt(8)
    body.text = lines[0]
    body.paragraphs[0].font.name = 'Calibri'
    body.paragraphs[0].font.size = Pt(20)
    body.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    body.paragraphs[0].space_after = Pt(12)

    for line in lines[1:]:
        p = body.add_paragraph()
        p.text = line
        p.level = 1
        p.font.name = 'Calibri'
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(8)


# Title slide
slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(slide_layout)
set_slide_background_white(slide)
slide.shapes.title.text = "ANCS Presentation"
slide.placeholders[1].text = (
    "عرض تقديمي يشرح نظام ANCS لتهيئة وإدارة الشبكات مع نقاط القوة، "
    "التكامل مع GNS3، وشاشات المشروع المهمة."
)
slide.placeholders[1].text_frame.paragraphs[0].font.name = 'Calibri'
slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(20)
slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

# Slide 2: What is ANCS?
blank_layout = presentation.slide_layouts[6]
slide = presentation.slides.add_slide(blank_layout)
set_slide_background_white(slide)
add_title(slide, "What is ANCS?")
add_paragraphs(slide, [
    "ANCS is an Autonomous Network Configuration & Orchestration System that helps engineers automate router and switch configuration, manage multi-vendor devices, and use AI guidance for network operations.",
    "The platform is designed for modern network labs, educational projects, and real-world multi-vendor environments.",
])

# Slide 3: The problem ANCS solves
slide = presentation.slides.add_slide(blank_layout)
set_slide_background_white(slide)
add_title(slide, "The problem ANCS solves")
add_paragraphs(slide, [
    "Traditional network configuration is slow, error-prone, and difficult to maintain across many devices.",
    "ANCS reduces manual mistakes by automating configuration generation, validation, and deployment.",
    "It also makes device management easier in GNS3 labs and multi-vendor networks.",
])

# Slide 4: Core features
slide = presentation.slides.add_slide(blank_layout)
set_slide_background_white(slide)
add_title(slide, "Core features")
add_paragraphs(slide, [
    "Guided Configuration Wizards for VLANs, OSPF, EIGRP, DHCP, ACLs, and more.",
    "Deep GNS3 integration for topology discovery and device control.",
    "AI Copilot for network discovery, troubleshooting, and auditing.",
    "Support for SSH, Telnet, multi-vendor devices, and bulk deployment.",
])

# Slide 5: Screenshots from the site
slide = presentation.slides.add_slide(blank_layout)
set_slide_background_white(slide)
add_title(slide, "Site Screenshots")
try:
    slide.shapes.add_picture(
        os.path.join(assets_dir, 'site-screenshot-hero.png'),
        Inches(0.5), Inches(1.6), width=Inches(4.5)
    )
    slide.shapes.add_picture(
        os.path.join(assets_dir, 'site-screenshot-preview.png'),
        Inches(5.1), Inches(1.6), width=Inches(4.5)
    )
except Exception:
    add_paragraphs(slide, [
        "Unable to load screenshot images. Please verify that the asset files exist at:",
        os.path.join(assets_dir, 'site-screenshot-hero.png'),
        os.path.join(assets_dir, 'site-screenshot-preview.png'),
    ])

# Slide 6: Final takeaways
slide = presentation.slides.add_slide(blank_layout)
set_slide_background_white(slide)
add_title(slide, "Final takeaways")
add_paragraphs(slide, [
    "ANCS makes network automation easier, faster, and more reliable.",
    "The platform is ideal for graduation project presentations and practical network lab demonstrations.",
    "Use this file as a complete overview of the website, the project goals, and the delivered features.",
])

presentation.save(output_path)
print(f"Created PowerPoint: {output_path}")
